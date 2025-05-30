import os
import json
import hashlib
from pathlib import Path
import pandas as pd
from pptx import Presentation
from docx import Document
import PyPDF2


class DataLoader:
    """
    DataLoader class to process various document formats and extract text chunks.
    It supports PDF, DOCX, PPTX, and CSV files, extracting text and metadata,
    and saving the processed chunks to a JSONL file.
    Attributes:
        corpus_dir (str): Directory containing the corpus files.
        processed_path (str): Path to the JSON file storing processed file hashes.
        output_path (str): Path to the output JSONL file for processed chunks.
        processed_hashes (set): Set of hashes for files that have already been processed.
    Methods:
        load_processed_hashes(): Loads processed file hashes from the JSON file.
        extract_text_from_pdf(file_path): Extracts text from a PDF file.
        extract_text_from_docx(file_path): Extracts text from a DOCX file.
        extract_text_from_pptx(file_path): Extracts text from a PPTX file.
        extract_text_from_csv(file_path): Extracts text from a CSV file.
        process_corpus(): Processes all files in the corpus directory, extracting text and metadata.
        update_processed_files(): Updates the JSON file with the current set of processed file hashes.
        save_chunks_as_jsonl(chunks): Saves the extracted text chunks to a JSONL file.
    """

    def __init__(
        self,
        corpus_dir="./corpus",
        processed_path="./processed_corpus/processed_files.json",
        output_path="./processed_corpus/processed_chunks.jsonl",
    ):
        self.corpus_dir = corpus_dir
        self.processed_path = processed_path
        self.output_path = output_path
        self.processed_hashes = self.load_processed_hashes()

    def load_processed_hashes(self):
        """
        Loads the set of processed file hashes from the JSON file.
        Returns:
            set: A set of file hashes that have already been processed.
        """
        if os.path.exists(self.processed_path):
            with open(self.processed_path, "r") as f:
                return set(json.load(f).get("file_hashes", []))
        return set()

    def extract_text_from_pdf(self, file_path):
        """
        Extracts text from a PDF file and returns it as a list of text chunks.
        Args:
            file_path (str): Path to the PDF file.
        Returns:
            list: A list of dictionaries, each containing 'content', 'type', and 'metadata'.
        """
        chunks = []
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text:
                    chunks.append(
                        {
                            "content": text.strip(),
                            "type": "text",
                            "metadata": {
                                "source": file_path,
                                "page_slide": f"Page {i+1}",
                                "embedded": False,
                            },
                        }
                    )
        return chunks

    def extract_text_from_docx(self, file_path):
        """
        Extracts text from a DOCX file and returns it as a list of text chunks.
        Args:
            file_path (str): Path to the DOCX file.
        Returns:
            list: A list of dictionaries, each containing 'content', 'type', and 'metadata'.
        """
        doc = Document(file_path)
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return [
            {
                "content": text.strip(),
                "type": "text",
                "metadata": {"source": file_path, "embedded": False},
            }
        ]

    def extract_text_from_pptx(self, file_path):
        """
        Extracts text from a PPTX file and returns it as a list of text chunks.
        Args:
            file_path (str): Path to the PPTX file.
        Returns:
            list: A list of dictionaries, each containing 'content', 'type', and 'metadata'.
        """
        prs = Presentation(file_path)
        chunks = []
        for i, slide in enumerate(prs.slides):
            texts = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
            if texts:
                chunks.append(
                    {
                        "content": "\n".join(texts).strip(),
                        "type": "text",
                        "metadata": {
                            "source": file_path,
                            "page_slide": f"Slide {i+1}",
                            "embedded": False,
                        },
                    }
                )
        return chunks

    def extract_text_from_csv(self, file_path):
        """
        Extracts text from a CSV file and returns it as a list of dictionaries representing the table.
        Args:
            file_path (str): Path to the CSV file.
        Returns:
            list: A list containing a single dictionary with 'content', 'type', and 'metadata'.
        """
        df = pd.read_csv(file_path)
        return [
            {
                "content": df.to_markdown(index=False),
                "type": "table",
                "metadata": {"source": file_path, "embedded": False},
            }
        ]

    def process_corpus(self):
        """
        Processes all files in the corpus directory, extracting text and metadata.
        Skips files that have already been processed based on their hash.
        Returns:
            list: A list of dictionaries, each containing 'content', 'type', and 'metadata'.
        """
        all_chunks = []
        for root, _, files in os.walk(self.corpus_dir):
            for file in files:
                file_path = os.path.join(root, file)
                file_hash = hashlib.sha256(file_path.encode()).hexdigest()

                if file_hash in self.processed_hashes:
                    print(f"Skipping already processed file: {file}")
                    continue

                ext = Path(file).suffix.lower()
                try:
                    if ext == ".pdf":
                        all_chunks.extend(self.extract_text_from_pdf(file_path))
                    elif ext == ".docx":
                        all_chunks.extend(self.extract_text_from_docx(file_path))
                    elif ext == ".pptx":
                        all_chunks.extend(self.extract_text_from_pptx(file_path))
                    elif ext == ".csv":
                        all_chunks.extend(self.extract_text_from_csv(file_path))
                    else:
                        print(f"Unsupported file format: {file}")
                        continue

                    self.processed_hashes.add(file_hash)

                except Exception as e:
                    print(f"Error processing {file_path}: {e}")
        return all_chunks

    def update_processed_files(self):
        """
        Updates the JSON file with the current set of processed file hashes.
        This ensures that files already processed are not reprocessed in future runs.
        """
        with open(self.processed_path, "w") as f:
            json.dump({"file_hashes": list(self.processed_hashes)}, f, indent=2)
        print(f"✅ Updated processed files: {len(self.processed_hashes)} total files.")

    def save_chunks_as_jsonl(self, chunks):
        """
        Saves the extracted text chunks to a JSONL file.
        Args:
            chunks (list): A list of dictionaries containing the extracted text chunks.
        """
        os.makedirs(os.path.dirname(self.output_path), exist_ok=True)
        with open(self.output_path, "w", encoding="utf-8") as f:
            for chunk in chunks:
                f.write(f"{chunk}\n")

    def run_pipeline(self):
        """
        Executes the data loading pipeline:
        1. Loads processed file hashes.
        2. Processes the corpus directory to extract text chunks.
        3. Saves the extracted chunks to a JSONL file.
        4. Updates the processed files JSON with new hashes.
        """
        chunks = self.process_corpus()
        if chunks:
            self.save_chunks_as_jsonl(chunks)
            self.update_processed_files()
            print(f"✅ Extracted {len(chunks)} chunks to {self.output_path}")
        else:
            print("📂 No new files found to process.")


if __name__ == "__main__":
    loader = DataLoader()
    loader.run_pipeline()
